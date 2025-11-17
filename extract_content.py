import os
import json
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_pptx_content(pptx_path, base_output_dir):
    """
    Extracts text and images from a PPTX file and saves them.

    Args:
        pptx_path (str): Path to the PPTX file.
        base_output_dir (str): The base directory for output (e.g., '00_Inbox').

    Returns:
        str: Path to the JSON file containing extracted data.
    """
    attachments_dir = os.path.join(base_output_dir, "attachments")
    if not os.path.exists(attachments_dir):
        os.makedirs(attachments_dir)

    prs = Presentation(pptx_path)
    
    extracted_data = {
        "slides": []
    }

    for i, slide in enumerate(prs.slides):
        slide_number = i + 1
        slide_content = {
            "slide_number": slide_number,
            "text": [],
            "images": []
        }

        # Extract text from shapes
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    slide_content["text"].append(text)

        # Extract images
        image_index = 0
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob
                image_ext = image.ext.lower()
                
                # Use a generic but unique name for now
                image_filename = f"slide_{slide_number}_image_{image_index + 1}.{image_ext}"
                image_path = os.path.join(attachments_dir, image_filename)
                
                with open(image_path, "wb") as f:
                    f.write(image_bytes)
                
                slide_content["images"].append(image_path)
                image_index += 1
        
        extracted_data["slides"].append(slide_content)

    # Save extracted data to a JSON file for the next steps
    json_output_path = os.path.join(base_output_dir, "extracted_content.json")
    with open(json_output_path, "w", encoding="utf-8") as f:
        json.dump(extracted_data, f, ensure_ascii=False, indent=4)

    print(f"Text and image metadata extracted to {json_output_path}")
    print(f"Images saved to {attachments_dir}")
    
    return json_output_path

if __name__ == "__main__":
    pptx_file = "00_Inbox/04_ARM汇编模拟器VisUAL.pptx"
    output_directory = "00_Inbox"
    extract_pptx_content(pptx_file, output_directory)
